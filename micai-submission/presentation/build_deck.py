#!/usr/bin/env python3
"""Build the CSHRL MICAI 2026 oral-presentation deck on the official template.

Usage: python3 build_deck.py
Reads  MICAI2026_Author_Template.pptx (official template, first slide mandatory)
Writes CSHRL_MICAI2026_oral.pptx (12 slides, paced for 10 min + 5 Q&A)
"""

import copy
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = "MICAI2026_Author_Template.pptx"
OUTPUT = "CSHRL_MICAI2026_oral.pptx"

# palette
BLUE = RGBColor(0x24, 0x56, 0xA6)      # MICAI logo blue
DARK = RGBColor(0x2E, 0x34, 0x40)      # body text
GRAY = RGBColor(0x5E, 0x68, 0x75)      # muted
RED = RGBColor(0xB3, 0x3A, 0x48)       # warnings / limits
GREEN = RGBColor(0x4F, 0x77, 0x3C)     # verified / strengths
CODE_BG = RGBColor(0x2E, 0x34, 0x40)   # code block background
CODE_FG = RGBColor(0xEC, 0xEF, 0xF4)
CODE_KW = RGBColor(0x88, 0xC0, 0xD0)
BOX_FILL = RGBColor(0xF0, 0xF4, 0xFA)

MONO = "Consolas"

prs = Presentation(TEMPLATE)
blank_layout = next(l for l in prs.slide_layouts if l.name == "Blank")


def set_runs(paragraph, text, size, color, bold=False, italic=False, font=None):
    """Fill a paragraph, honoring **bold** markers inside `text`."""
    for i, part in enumerate(re.split(r"\*\*", text)):
        if part == "":
            continue
        run = paragraph.add_run()
        run.text = part
        f = run.font
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold or (i % 2 == 1)
        f.italic = italic
        if font:
            f.name = font


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.30), Inches(12.2), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    set_runs(tf.paragraphs[0], text, 30, BLUE, bold=True)
    return box


def add_body(slide, items, left=0.55, top=1.35, width=12.2, height=5.2):
    """items: list of (text, kind) where kind in bullet|sub|formula|muted."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, kind in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if kind == "bullet":
            set_runs(p, "•  " + text, 19, DARK)
            p.space_after = Pt(10)
        elif kind == "sub":
            set_runs(p, "–  " + text, 16, GRAY)
            p.level = 1
            p.space_after = Pt(6)
        elif kind == "formula":
            set_runs(p, text, 21, BLUE, italic=True)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(8)
            p.space_after = Pt(12)
        elif kind == "muted":
            set_runs(p, text, 14, GRAY, italic=True)
            p.space_before = Pt(6)
    return box


def add_box(slide, text, left, top, width, height, accent, size=16):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BOX_FILL
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.10)
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        set_runs(p, line, size, DARK)
        p.space_after = Pt(4)
    return shape


def new_slide():
    return prs.slides.add_slide(blank_layout)


# ---------------------------------------------------------------- cover (slide 0)
def fill_placeholder(tf, lines):
    """Replace a placeholder's paragraphs with (text, size) lines, keeping theme colors."""
    tf.clear()
    for i, (text, size) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)

cover = prs.slides[0]
for shape in cover.shapes:
    if not shape.is_placeholder:
        continue
    if shape.placeholder_format.type == 3:  # CENTER_TITLE
        fill_placeholder(shape.text_frame, [
            ("Coinductive Symmetric Homomorphism Reinforcement Learning", 30),
            ("Optimality as Structure Preservation", 22),
        ])
    elif shape.placeholder_format.type == 4:  # SUBTITLE
        fill_placeholder(shape.text_frame, [
            ("Ricardo Corral-Corral  ·  ricardo@gyx.ai", 18),
        ])

# NOTE: the 7 template example slides are removed at the END of this script,
# after the content slides are appended, so python-pptx never reuses the freed
# relationship ids (Keynote's importer chokes on reused/shuffled rIds).

# ---------------------------------------------------------------- slide 2
s = new_slide()
add_title(s, "The problem with scalar returns")
add_body(s, [
    ("G  =  r₀ + γ·r₁ + γ²·r₂ + ⋯        γ ∈ [0, 1)", "formula"),
    ("The infinite future is **collapsed into one number** — all temporal structure is discarded.", "bullet"),
    ("γ is **arbitrary**: it makes the sum converge; the task does not demand it.", "bullet"),
    ("A reward 10 steps ahead counts for γ¹⁰ of its face value — the math shapes the objective.", "bullet"),
], height=3.4)
add_box(s, "Consequence: in **sacrifice patterns** — pay now, gain forever — scalar methods can provably pick the wrong action.",
        0.55, 5.3, 12.2, 1.0, RED, size=17)

# ---------------------------------------------------------------- slide 3
s = new_slide()
add_title(s, "Key idea: keep the stream")
add_body(s, [
    ("From any state, record the **best achievable reward at every future depth** — never collapse it:", "bullet"),
    ("value(s′)  =  [ v₀, v₁, v₂, … ]    — an infinite stream, not a number", "formula"),
    ("Compare states by **pointwise stream dominance**:", "bullet"),
    ("x ⊑ y   ⟺   x(t) ≤ y(t)   for every t = 0, 1, 2, …", "formula"),
    ("No discount factor.  No averaging.  No horizon cutoff.", "bullet"),
    ("The better state must be at least as good at **every single future depth**.", "bullet"),
])

# ---------------------------------------------------------------- slide 4
s = new_slide()
add_title(s, "Two coinductive optimality conditions")
add_box(s, "**CoinductiveHomomorphism — strategic**\n"
           "a ≼ b  ⟹  value(next(s,a)) ⊑ value(next(s,b))\n"
           "Judge actions by **where they lead**. Immediate reward is invisible; successor quality is everything.",
        0.55, 1.45, 5.95, 2.5, GREEN)
add_box(s, "**CoindHomo — tactical**\n"
           "a ≼ b  ⟹  action-value(s,a) ⊑ action-value(s,b)\n"
           "The better action must win the **immediate payoff and the long-term future** — head and tail.",
        6.80, 1.45, 5.95, 2.5, BLUE)
add_box(s, "**Decomposition theorem (machine-checked):**  CoindHomo  =  CoinductiveHomomorphism  +  immediate-reward compatibility.",
        0.55, 4.35, 12.2, 1.0, BLUE, size=17)

# ---------------------------------------------------------------- slide 5
s = new_slide()
add_title(s, "Why “Symmetric”: the full ranking")
add_body(s, [
    ("The object of interest is the **entire order** on the action set — a permutation, an element of the **symmetric group** — not just the argmax.", "bullet"),
    ("An agent that has understood its environment can say **which of its two worst options is less bad**. An argmax agent cannot.", "bullet"),
    ("A top-action-only learner can look competent while its model is **structurally faulty** — until its preferred action disappears.", "bullet"),
], height=3.3)
add_box(s, "The full-ranking stance pays operationally:\n"
           "•  learning converges by fixing **every** violated pair — no exploration mechanism needed\n"
           "•  top action unavailable ⟹ **O(1)** fallback, with a verified guarantee",
        0.55, 4.9, 12.2, 1.6, GREEN)

# ---------------------------------------------------------------- slide 6
s = new_slide()
add_title(s, "Verified counterexamples: the strict gap")
add_body(s, [
    ("Three machine-checked environments separate the two conditions:", "bullet"),
    ("**BinarySacrifice** — trap pays 1 now; paradise pays forever", "sub"),
    ("**SkillInvestment** — train (0 now) vs. work (2 now)", "sub"),
    ("**PreparationDilemma** — prepare vs. act immediately", "sub"),
    ("Each has a CoinductiveHomomorphism proof + a proof that **no CoindHomo can rank** the sacrificing actions.", "bullet"),
], height=3.3)
add_box(s, "**Q-learning fails here, provably:** for γ < ½ it selects the action with higher immediate reward and the strictly inferior successor — the discount geometrically suppresses the evidence.  (QLearningFailure.agda)",
        0.55, 4.9, 12.2, 1.35, RED, size=16)

# ---------------------------------------------------------------- slide 7
s = new_slide()
add_title(s, "The Agda core")
code = slide_code = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.45), Inches(7.1), Inches(2.9))
code.fill.solid()
code.fill.fore_color.rgb = CODE_BG
code.line.fill.background()
tf = code.text_frame
tf.word_wrap = False
tf.margin_left = Inches(0.25)
tf.margin_top = Inches(0.15)
code_lines = [
    ("-- stream dominance: an infinite witness", GRAY),
    ("record _⊑_ (x y : Stream R) : Set where", CODE_FG),
    ("  coinductive", CODE_KW),
    ("  field", CODE_KW),
    ("    head≤ : head x ≤ head y", CODE_FG),
    ("    tail⊑ : tail x ⊑ tail y", CODE_FG),
]
first = True
for line, color in code_lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run()
    run.text = line
    run.font.name = MONO
    run.font.size = Pt(16)
    run.font.color.rgb = color
caption = s.shapes.add_textbox(Inches(0.55), Inches(4.45), Inches(7.1), Inches(0.4))
set_runs(caption.text_frame.paragraphs[0],
         "(notation lightly simplified; the paper writes _≤ₛ_)", 12, GRAY, italic=True)
add_body(s, [
    ("Compiles under **--safe --guardedness**: zero postulates, no escape hatches.", "bullet"),
    ("Copatterns make coinductive proofs direct-style — why **Agda** over Coq / Isabelle / Lean.", "bullet"),
    ("The paper is literate Agda: **building the PDF type-checks every definition shown**.", "bullet"),
], left=7.95, top=1.45, width=4.8, height=4.5)

# ---------------------------------------------------------------- slide 8
s = new_slide()
add_title(s, "From ideal to algorithm: the Finder")
add_body(s, [
    ("The full coinductive condition is **undecidable in general** — like AIXI, an ideal organizing practical approximations.", "bullet"),
    ("The Finder compares **finite traces lexicographically** at depth n: earlier depths dominate, ties recurse deeper.", "bullet"),
    ("**No discount factor required** — the lexicographic order is well-defined at any depth.", "bullet"),
], height=3.3)
add_box(s, "**Exactness is not assumed — it is verified.**  Environment Classes prove conditions (e.g. horizon sufficiency) under which finite traces decide the infinite condition.",
        0.55, 4.9, 12.2, 1.2, GREEN, size=17)

# ---------------------------------------------------------------- slide 9
s = new_slide()
add_title(s, "Learning as symmetry restoration")
add_body(s, [
    ("A **violation**: the ranking disagrees with the ground-truth comparator on one pair, at one state.", "bullet"),
    ("The repair is a **transposition** — and transpositions generate the symmetric group: learning literally walks it.", "bullet"),
    ("max-violations  ≤  |S| × |A|(|A|−1) / 2", "formula"),
], height=3.0)
add_box(s, "**Swap Monotonicity (verified):** each swap fixes its pair, leaves unrelated pairs unchanged, and strictly decreases total violations  ⟹  convergence in bounded steps.",
        0.55, 4.45, 12.2, 1.1, GREEN, size=16)
foot = s.shapes.add_textbox(Inches(0.55), Inches(5.75), Inches(12.2), Inches(0.7))
set_runs(foot.text_frame.paragraphs[0],
         "Corollaries: uniform pairwise sampling suffices (exploration/exploitation dissolves) · O(1) adaptation via demote-preserves-dominance.",
         13, GRAY, italic=True)

# ---------------------------------------------------------------- slide 10
s = new_slide()
add_title(s, "Stochastic extension")
add_body(s, [
    ("Environments lift through the **finite distribution monad**; values become expected-reward streams.", "bullet"),
    ("Pointwise order is too strict under expectation ⟹ **lexicographic coinductive comparison**: compare heads, recurse on ties.", "bullet"),
    ("SD[0] = FOSD   ⊂   SD[1] = SOSD   ⊂   SD[2]   ⊂   ⋯", "formula"),
    ("Each level captures a broader class of utility functions — all verifiable.", "bullet"),
    ("Verified rankings **compose**: product, convolution, sum, scaling — zero extra proof burden for independent subsystems.", "bullet"),
])

# ---------------------------------------------------------------- slide 11
s = new_slide()
add_title(s, "Scaling verification: Environment Classes")
add_body(s, [
    ("Reusable verification infrastructure between core theory and tasks:", "bullet"),
    ("**FiniteDeterministicMDP** — TwoState, GridWorld, KeyTreasure, the sacrifice tasks…", "sub"),
    ("**CombinatorialPlacementMDP** — Queens1, Queens8", "sub"),
    ("**StochasticFiniteMDP** — CoinFlip, GamblersRuin, RandomWalk, BiasedBandit", "sub"),
], height=3.0)
add_box(s, "**8-Queens: full CoindHomo over 8⁸ ≈ 16.7 million candidate placements — zero postulates.**\n"
           "The EC supplies the finite-trace-to-infinite-stream bridge; the task author supplies four simple ingredients.",
        0.55, 4.55, 12.2, 1.5, GREEN, size=17)

# ---------------------------------------------------------------- slide 12
s = new_slide()
add_title(s, "Takeaways")
add_body(s, [
    ("**Optimality = structure preservation**, not scalar maximization.", "bullet"),
    ("The **full ranking** is the object of knowledge — and it pays: provable convergence, instant adaptation.", "bullet"),
    ("**Verified ≠ small**: environment classes carry the proofs to 16.7M-configuration domains.", "bullet"),
    ("Honest limits: finite action spaces; exact traces need the step function (model-based flavor); empirical evaluation ongoing.", "muted"),
], height=3.4)
add_box(s, "**github.com/doctorcorral/csh-rl**\ncomplete Agda artifact · the paper type-checks itself\n\n**Thank you — questions?**",
        3.1, 4.8, 7.1, 1.7, BLUE, size=18)

# ------------------------------------------------- drop the 7 example slides
# (done last: freed rIds are not reused, keeping the rel map sequential)
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst)[1:8]:
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)

prs.save(OUTPUT)
print(f"wrote {OUTPUT}: {len(prs.slides._sldIdLst)} slides")
